import os
import sys

def generate_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx is not installed. Please install it with 'pip install python-pptx'")
        sys.exit(1)

    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "OOP-Verwaltung in Python: Ticket-System"
    subtitle.text = "Ein Abschluss-Projekt von Bennet, Timothy & Jason\nKlasse: 10its2 | LF05"

    # Slide 2: Projektziel
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "1. Projektziel"
    tf = body_shape.text_frame
    tf.text = "Entwicklung eines Ticket-Management-Systems"
    p = tf.add_paragraph()
    p.text = "Strikte Anwendung von Objektorientierung (OOP)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Persistente Speicherung in CSV"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zielgruppe: IT-Abteilungen für Support-Anfragen"
    p.level = 1

    # Slide 3: Funktionalität (Features)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "2. Funktionalität & Features"
    tf = body_shape.text_frame
    tf.text = "CRUD-Operationen (Create, Read, Update, Delete)"
    p = tf.add_paragraph()
    p.text = "Grafische Oberfläche mit Tkinter (professionelles Design)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zwei Logik-Methoden: Suche und Prioritäten-Statistik"
    p.level = 1

    # Slide 4: Architektur (3-Schichten-Modell)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "3. Architektur (3-Schichten-Modell)"
    tf = body_shape.text_frame
    tf.text = "Das System ist strikt getrennt:"
    p = tf.add_paragraph()
    p.text = "Fachklasse (Ticket): Hält nur Daten und Getter."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Datenhaltung (FileManager): Liest und schreibt CSV-Dateien."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Benutzeroberfläche (TicketSystem): Verbindet GUI mit Logik."
    p.level = 1

    # Slide 5: OOP Leitfragen - Kapselung
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "4. OOP-Konzepte: Datenkapselung"
    tf = body_shape.text_frame
    tf.text = "Schutz der Daten vor unerlaubtem Zugriff"
    p = tf.add_paragraph()
    p.text = "Attribute sind private (z.B. __ticket_id)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zugriff nur über Getter-Methoden (get_title())"
    p.level = 1

    # Slide 6: Code-Snippet (Klasse Ticket)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "5. Code-Snippet: Kapselung"
    tf = body_shape.text_frame
    tf.text = "class Ticket:\n    def __init__(self, ticket_id, title, priority, text):\n        self.__ticket_id = ticket_id\n        self.__title = title\n\n    def get_title(self):\n        return self.__title"
    tf.paragraphs[0].font.size = Pt(16)

    # Slide 7: UML Diagramme (Übersicht)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "6. UML Klassendiagramm"
    tf = body_shape.text_frame
    tf.text = "TicketSystem (UI)"
    p = tf.add_paragraph()
    p.text = "|--> FileManager (Daten)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "       |--> Ticket (Fachklasse)"
    p.level = 2

    # Slide 8: Reflexion & Fazit
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "7. Fazit & Reflexion"
    tf = body_shape.text_frame
    tf.text = "Was lief gut?"
    p = tf.add_paragraph()
    p.text = "Saubere Trennung der Schichten"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Erfolgreiche Umsetzung der strikten Vorgaben"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Was haben wir gelernt?"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Wichtigkeit von Interfaces und Namenskonventionen"
    p.level = 1

    prs.save("Praesentation.pptx")
    print("Presentation created successfully as 'Praesentation.pptx'")

if __name__ == "__main__":
    generate_pptx()
